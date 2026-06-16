"""OpenAlio MCP — McKinsey/BCG 기업 컨설팅 스타일"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W = Inches(13.33)
H = Inches(7.5)

# ── 색상 시스템 (정제된 기업 팔레트) ────────────────────────────
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BG     = RGBColor(0xF8, 0xFA, 0xFD)   # 슬라이드 배경

NAVY_D = RGBColor(0x0C, 0x2A, 0x54)   # 헤딩 / 표지
NAVY   = RGBColor(0x1A, 0x47, 0x89)   # 주 액센트
NAVY_L = RGBColor(0xCC, 0xDB, 0xF4)   # 틴트

TEAL   = RGBColor(0x00, 0x7A, 0x73)   # 2차 액센트
TEAL_L = RGBColor(0xC8, 0xE8, 0xE6)

GOLD   = RGBColor(0xBF, 0x7C, 0x1C)   # 강조 (Key callout)
GOLD_L = RGBColor(0xFC, 0xED, 0xCF)

RUST   = RGBColor(0xB3, 0x47, 0x28)   # Before / 문제
RUST_L = RGBColor(0xF8, 0xE1, 0xD8)

PLUM   = RGBColor(0x50, 0x37, 0x87)   # 보조
PLUM_L = RGBColor(0xE4, 0xDE, 0xF8)

SAGE   = RGBColor(0x24, 0x6C, 0x3D)   # After / 해결
SAGE_L = RGBColor(0xD0, 0xEC, 0xDB)

TXT_D  = RGBColor(0x0F, 0x18, 0x28)
TXT_M  = RGBColor(0x37, 0x45, 0x5D)
TXT_L  = RGBColor(0x5F, 0x6F, 0x89)
TXT_XL = RGBColor(0x95, 0xA5, 0xB9)

BORDER = RGBColor(0xD2, 0xDC, 0xEC)
LINE   = RGBColor(0xBD, 0xCB, 0xE1)
CARD   = RGBColor(0xF0, 0xF4, 0xF9)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── 레이아웃 상수 ─────────────────────────────────────────────
LM   = Inches(0.42)
CW   = W - LM * 2
C2W  = (CW - Inches(0.4)) / 2
C2X2 = LM + C2W + Inches(0.4)
C3W  = (CW - Inches(0.28) * 2) / 3
C3X2 = LM + C3W + Inches(0.28)
C3X3 = LM + (C3W + Inches(0.28)) * 2

# ── 헬퍼 ─────────────────────────────────────────────────────
def R(sl, x, y, w, h, fill=None, lc=None, lw=Pt(0.75)):
    s = sl.shapes.add_shape(1, x, y, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else:    s.fill.background()
    s.line.fill.background()
    if lc:   s.line.color.rgb = lc; s.line.width = lw
    return s

def T(sl, text, x, y, w, h, sz=13, bold=False,
      color=TXT_M, align=PP_ALIGN.LEFT, italic=False):
    b = sl.shapes.add_textbox(x, y, w, h)
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return b

def AP(tf, text, sz=13, bold=False, color=TXT_M,
       align=PP_ALIGN.LEFT, sp=Pt(6)):
    p = tf.add_paragraph(); p.alignment = align; p.space_before = sp
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = color

def new_slide(num, title, sub=None, total=15):
    """표준 헤더 슬라이드"""
    sl = prs.slides.add_slide(BLANK)
    R(sl, 0, 0, W, H, fill=BG)
    R(sl, 0, 0, W, Inches(0.04), fill=NAVY)
    R(sl, 0, Inches(0.04), Inches(0.055), Inches(1.1), fill=NAVY)
    T(sl, title, LM, Inches(0.06), Inches(11.2), Inches(0.62),
      sz=22, bold=True, color=NAVY_D)
    T(sl, f"{num}  /  {total}", W - Inches(1.5), Inches(0.06),
      Inches(1.3), Inches(0.5), sz=10, color=TXT_XL, align=PP_ALIGN.RIGHT)
    cy = Inches(0.68)
    if sub:
        T(sl, sub, LM, cy, Inches(11.2), Inches(0.36),
          sz=12, color=TXT_L, italic=True)
        cy += Inches(0.4)
    R(sl, LM, cy + Inches(0.04), CW, Inches(0.01), fill=LINE)
    R(sl, 0, H - Inches(0.26), W, Inches(0.26), fill=NAVY_D)
    T(sl, "OpenAlio MCP  ·  재경부 공공정책국",
      LM, H - Inches(0.24), Inches(7), Inches(0.2), sz=8, color=NAVY_L)
    return sl, cy + Inches(0.18)

def acard(sl, x, y, w, h, accent=NAVY, bg=WHITE):
    """왼쪽 액센트 바 + 테두리 카드"""
    R(sl, x, y, w, h, fill=bg, lc=BORDER, lw=Pt(0.5))
    R(sl, x, y, Inches(0.048), h, fill=accent)

def ctitle(sl, text, x, y, w, accent=NAVY):
    """카드 섹션 제목"""
    T(sl, text, x + Inches(0.14), y, w - Inches(0.2), Inches(0.44),
      sz=14, bold=True, color=accent)

def blist(sl, items, x, y, w, h, sz=13, color=TXT_M, sp=Pt(7)):
    bx = sl.shapes.add_textbox(x, y, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        if first: p = tf.paragraphs[0]; first = False
        else:     p = tf.add_paragraph(); p.space_before = sp
        r = p.add_run(); r.text = f"·  {item}"
        r.font.size = Pt(sz); r.font.color.rgb = color

def stepbox(sl, x, y, w, h, n, label, accent=NAVY):
    R(sl, x, y, w, h, fill=WHITE, lc=BORDER, lw=Pt(0.5))
    R(sl, x, y, Inches(0.36), h, fill=accent)
    T(sl, str(n), x, y, Inches(0.36), h, sz=13, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)
    T(sl, label, x + Inches(0.44), y + Inches(0.06),
      w - Inches(0.52), h - Inches(0.12), sz=13, color=TXT_M)

def arr_dn(sl, x, y, h=Inches(0.2)):
    R(sl, x - Inches(0.025), y, Inches(0.05), h, fill=TXT_XL)

def hbar(sl, text, accent=NAVY, y=H - Inches(0.75)):
    R(sl, LM, y, CW, Inches(0.52), fill=accent)
    T(sl, text, LM + Inches(0.2), y + Inches(0.08),
      CW - Inches(0.4), Inches(0.36), sz=14, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)

def tag(sl, text, x, y, accent=NAVY):
    R(sl, x, y, Inches(0.07), Inches(0.07), fill=accent)
    T(sl, text, x + Inches(0.15), y - Inches(0.04),
      Inches(5), Inches(0.28), sz=11, bold=True, color=accent)


# ════════════════════════════════════════════════════════════════
#  S00 — 표지
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
# 왼쪽 패널
R(sl, 0, 0, Inches(6.9), H, fill=NAVY_D)
R(sl, 0, 0, Inches(6.9), Inches(0.06), fill=TEAL)
# 오른쪽 패널
R(sl, Inches(6.9), 0, W - Inches(6.9), H, fill=BG)
# 오른쪽 장식 (절제된 기하학)
R(sl, Inches(7.2), Inches(0.5), Inches(5.7), Inches(0.01), fill=LINE)
R(sl, Inches(7.2), Inches(1.2), Inches(5.7), Inches(0.01), fill=LINE)
R(sl, Inches(7.2), Inches(1.9), Inches(5.7), Inches(0.01), fill=LINE)
R(sl, Inches(7.5), Inches(0.5), Inches(0.01), Inches(6.5), fill=LINE)
R(sl, Inches(7.5), Inches(0.5), Inches(0.25), Inches(0.25), fill=NAVY_L)
R(sl, Inches(8.2), Inches(1.5), Inches(4.5), Inches(3.2), fill=CARD,
  lc=BORDER, lw=Pt(0.5))
R(sl, Inches(8.2), Inches(1.5), Inches(0.048), Inches(3.2), fill=TEAL)
T(sl, "핵심 데이터",
  Inches(8.45), Inches(1.65), Inches(4.1), Inches(0.42),
  sz=12, bold=True, color=TEAL)
stats = [
    ("355개", "공공기관"),
    ("11개", "메트릭 카테고리"),
    ("32개", "AI 도구"),
    ("6년치", "시계열 데이터"),
]
sy = Inches(2.2)
for val, lbl in stats:
    T(sl, val, Inches(8.45), sy, Inches(2.1), Inches(0.45),
      sz=20, bold=True, color=NAVY_D)
    T(sl, lbl, Inches(10.6), sy + Inches(0.06), Inches(1.8), Inches(0.35),
      sz=12, color=TXT_L)
    sy += Inches(0.58)
# 왼쪽 콘텐츠
T(sl, "공공기관 경영정보 AI 활용 실험",
  Inches(0.55), Inches(0.7), Inches(6.0), Inches(0.38),
  sz=11, color=TEAL_L, italic=True)
T(sl, "OpenAlio MCP",
  Inches(0.55), Inches(1.2), Inches(6.0), Inches(1.5),
  sz=46, bold=True, color=WHITE)
R(sl, Inches(0.55), Inches(2.82), Inches(3.8), Inches(0.055), fill=TEAL)
T(sl, "공공기관 정보공개의 AI 활용성 제고를 위한",
  Inches(0.55), Inches(3.0), Inches(6.0), Inches(0.42),
  sz=15, color=NAVY_L)
T(sl, "MCP 기반 실험",
  Inches(0.55), Inches(3.42), Inches(6.0), Inches(0.42),
  sz=15, color=NAVY_L)
R(sl, Inches(0.55), Inches(5.55), Inches(4.2), Inches(0.01), fill=LINE)
T(sl, "김보성 사무관",
  Inches(0.55), Inches(5.68), Inches(5.5), Inches(0.5),
  sz=18, bold=True, color=WHITE)
T(sl, "재경부 공공정책국  |  2026. 06.",
  Inches(0.55), Inches(6.2), Inches(5.5), Inches(0.36),
  sz=13, color=NAVY_L)


# ════════════════════════════════════════════════════════════════
#  S01 — ALIO 소개
# ════════════════════════════════════════════════════════════════
sl, cy = new_slide(1, "01   공공기관 정보공개의 성공, ALIO",
                   "ALIO는 공공기관 투명성 제고의 대표적 성공 사례")

# 왼쪽: 핵심 성과
acard(sl, LM, cy, C2W, Inches(5.6))
ctitle(sl, "ALIO 핵심 성과", LM, cy + Inches(0.1), C2W)
R(sl, LM + Inches(0.14), cy + Inches(0.56), C2W - Inches(0.28),
  Inches(0.01), fill=LINE)

feats = [
    ("공공기관 경영정보 대국민 공개", "재무·인사·경영평가 등 50개 공시항목"),
    ("국민 알권리 확대", "누구나 무료로 열람 가능한 개방형 정보"),
    ("공공기관 책임성 강화", "공시 의무화로 투명성·책임성 제고"),
    ("지속적 서비스 확장", "355개 공공기관 정보 통합 제공"),
]
fy = cy + Inches(0.72)
for main, desc in feats:
    R(sl, LM + Inches(0.14), fy, Inches(0.07), Inches(0.07), fill=NAVY)
    T(sl, main, LM + Inches(0.3), fy - Inches(0.03),
      C2W - Inches(0.44), Inches(0.34), sz=13, bold=True, color=TXT_D)
    T(sl, desc, LM + Inches(0.3), fy + Inches(0.3),
      C2W - Inches(0.44), Inches(0.32), sz=12, color=TXT_L)
    fy += Inches(1.05)

# 오른쪽: 발전 연혁
acard(sl, C2X2, cy, C2W, Inches(5.6), accent=TEAL)
ctitle(sl, "서비스 발전 연혁", C2X2, cy + Inches(0.1), C2W, accent=TEAL)
R(sl, C2X2 + Inches(0.14), cy + Inches(0.56), C2W - Inches(0.28),
  Inches(0.01), fill=LINE)

timeline = [
    ("2005", "ALIO 출범", "공공기관 경영정보 최초 온라인 공개 시작", NAVY),
    ("2014", "ALIO PLUS", "채용·시설·사업 등 생활밀착 정보 통합", TEAL),
    ("2019", "JOB-ALIO", "공공기관 채용 전문 플랫폼 분리 출범", PLUM),
    ("2026", "OpenAlio MCP", "AI 활용 가능 개방형 인터페이스 실험", GOLD),
]
ty = cy + Inches(0.72)
for year, name, desc, col in timeline:
    R(sl, C2X2 + Inches(0.14), ty, Inches(0.72), Inches(0.78), fill=col)
    T(sl, year, C2X2 + Inches(0.14), ty + Inches(0.18),
      Inches(0.72), Inches(0.4), sz=12, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)
    T(sl, name, C2X2 + Inches(1.0), ty + Inches(0.02),
      C2W - Inches(1.2), Inches(0.36), sz=14, bold=True, color=col)
    T(sl, desc, C2X2 + Inches(1.0), ty + Inches(0.4),
      C2W - Inches(1.2), Inches(0.32), sz=12, color=TXT_L)
    if year != "2026":
        R(sl, C2X2 + Inches(0.49), ty + Inches(0.78),
          Inches(0.025), Inches(0.22), fill=TXT_XL)
    ty += Inches(1.05)


# ════════════════════════════════════════════════════════════════
#  S02 — AI 시대의 과제
# ════════════════════════════════════════════════════════════════
sl, cy = new_slide(2, "02   그러나 AI 시대의 새로운 과제",
                   "정보 공개는 성공했지만 정보 활용은 여전히 어렵다")

acard(sl, LM, cy, C2W, Inches(5.55))
ctitle(sl, "분산된 공공기관 정보원", LM, cy + Inches(0.1), C2W)
R(sl, LM + Inches(0.14), cy + Inches(0.54),
  C2W - Inches(0.28), Inches(0.01), fill=LINE)

srcs = [
    ("ALIO", "경영공시 정보 (재무·인사·경영평가)", NAVY),
    ("ALIO PLUS", "채용·시설·사업 정보 통합", TEAL),
    ("JOB-ALIO", "채용 공고 전문 서비스", PLUM),
    ("기관 홈페이지", "기관별 자체 정보·공고", RUST),
    ("법령 / 지침", "법적 근거·운영 지침", SAGE),
    ("뉴스 / 보도자료", "실시간 이슈 및 보도", GOLD),
]
sy2 = cy + Inches(0.7)
for name, desc, col in srcs:
    R(sl, LM + Inches(0.14), sy2, C2W - Inches(0.28),
      Inches(0.62), fill=WHITE, lc=BORDER, lw=Pt(0.5))
    R(sl, LM + Inches(0.14), sy2, Inches(0.04),
      Inches(0.62), fill=col)
    T(sl, name, LM + Inches(0.3), sy2 + Inches(0.06),
      Inches(1.8), Inches(0.3), sz=13, bold=True, color=col)
    T(sl, desc, LM + Inches(2.2), sy2 + Inches(0.15),
      C2W - Inches(2.4), Inches(0.3), sz=12, color=TXT_L)
    sy2 += Inches(0.73)

acard(sl, C2X2, cy, C2W, Inches(5.55), accent=RUST)
ctitle(sl, "AI 활용이 어려운 이유", C2X2, cy + Inches(0.1), C2W, accent=RUST)
R(sl, C2X2 + Inches(0.14), cy + Inches(0.54),
  C2W - Inches(0.28), Inches(0.01), fill=LINE)

problems = [
    ("분산성", "여러 사이트에 나뉜 정보 — 통합 조회 불가"),
    ("비구조성", "PDF·엑셀·웹페이지 혼재 — AI 파싱 어려움"),
    ("비연결성", "기관 간 연계 없어 비교·분석 불가"),
    ("저기계가독성", "표준화된 API 부재 — AI 직접 활용 불가"),
]
py = cy + Inches(0.72)
for i, (ptitle, pdesc) in enumerate(problems):
    R(sl, C2X2 + Inches(0.14), py, C2W - Inches(0.28),
      Inches(1.1), fill=RUST_L if i % 2 == 0 else WHITE,
      lc=BORDER, lw=Pt(0.5))
    T(sl, ptitle, C2X2 + Inches(0.28), py + Inches(0.1),
      C2W - Inches(0.5), Inches(0.36), sz=14, bold=True, color=RUST)
    T(sl, pdesc, C2X2 + Inches(0.28), py + Inches(0.48),
      C2W - Inches(0.5), Inches(0.52), sz=12, color=TXT_M)
    py += Inches(1.22)

T(sl, "정보는 공개됐지만, AI는 아직 '읽지 못한다'",
  C2X2 + Inches(0.14), cy + Inches(5.24), C2W - Inches(0.28),
  Inches(0.32), sz=12, bold=True, color=RUST, italic=True)


# ════════════════════════════════════════════════════════════════
#  S03 — 국민 Pain Point
# ════════════════════════════════════════════════════════════════
sl, cy = new_slide(3, "03   국민 입장에서의 Pain Point",
                   "정보는 공개되어 있으나 '질문하기' 어렵다")

# 질문 강조
R(sl, LM, cy, CW, Inches(0.62), fill=NAVY_D)
R(sl, LM, cy, Inches(0.055), Inches(0.62), fill=GOLD)
T(sl, '  "한국전력의 최근 5년 부채 추이는?"',
  LM + Inches(0.2), cy + Inches(0.1), CW - Inches(0.4), Inches(0.44),
  sz=18, bold=True, color=WHITE)
cy += Inches(0.74)

# Before
acard(sl, LM, cy, C2W, Inches(4.95), accent=RUST)
ctitle(sl, "Before  —  현재 절차", LM, cy + Inches(0.1), C2W, accent=RUST)
R(sl, LM + Inches(0.14), cy + Inches(0.54),
  C2W - Inches(0.28), Inches(0.01), fill=LINE)
bsteps = ["ALIO 접속", "기관 검색", "재무현황 이동",
          "자료 다운로드", "직접 분석"]
sy3 = cy + Inches(0.67)
for i, s in enumerate(bsteps, 1):
    stepbox(sl, LM + Inches(0.14), sy3, C2W - Inches(0.28),
            Inches(0.52), i, s, RUST)
    if i < len(bsteps):
        arr_dn(sl, LM + Inches(0.14) + (C2W - Inches(0.28)) / 2,
               sy3 + Inches(0.52), Inches(0.18))
    sy3 += Inches(0.7)
T(sl, "⏱  수십 분 ~ 수 시간 소요",
  LM + Inches(0.14), cy + Inches(4.72), C2W - Inches(0.28), Inches(0.28),
  sz=11, bold=True, color=RUST, italic=True)

# After
acard(sl, C2X2, cy, C2W, Inches(4.95), accent=SAGE)
ctitle(sl, "After  —  OpenAlio MCP 활용", C2X2, cy + Inches(0.1), C2W, accent=SAGE)
R(sl, C2X2 + Inches(0.14), cy + Inches(0.54),
  C2W - Inches(0.28), Inches(0.01), fill=LINE)

asteps = [
    ("1", "자연어로 질문 입력"),
    ("2", "AI가 MCP 도구 자동 선택·호출"),
    ("3", "ALIO 데이터 실시간 조회"),
    ("4", "분석 · 시각화 수행"),
    ("5", "결과 + 출처 · 공시일 반환"),
]
sy4 = cy + Inches(0.67)
for n, s in asteps:
    stepbox(sl, C2X2 + Inches(0.14), sy4, C2W - Inches(0.28),
            Inches(0.52), n, s, SAGE)
    if n != "5":
        arr_dn(sl, C2X2 + Inches(0.14) + (C2W - Inches(0.28)) / 2,
               sy4 + Inches(0.52), Inches(0.18))
    sy4 += Inches(0.7)
T(sl, "⚡  수 초 이내 완료 · 출처 자동 첨부",
  C2X2 + Inches(0.14), cy + Inches(4.72), C2W - Inches(0.28), Inches(0.28),
  sz=11, bold=True, color=SAGE, italic=True)


# ════════════════════════════════════════════════════════════════
#  S04 — 공공기관 실무자 Pain Point
# ════════════════════════════════════════════════════════════════
sl, cy = new_slide(4, "04   공공기관 실무자 입장에서의 Pain Point",
                   "벤치마킹과 업무학습 비용이 높다")

pain_cards = [
    ("유사기관 운영사례 탐색",
     "동종 기관 운영사례를 여러 사이트에서\n수작업으로 수집해야 함",
     "AI에게 물어보면 즉시 정리된\n벤치마킹 결과 제공",
     NAVY),
    ("경영평가 대응사례 확인",
     "지적사항·우수사례 PDF를 연도별로\n일일이 찾아서 검토",
     "편람 키워드 검색 + 연도별\n지적 패턴 즉시 분석",
     TEAL),
    ("신규 담당자 업무학습",
     "법령·지침·편람이 분산돼\n온보딩 비용 과다",
     "법령·지침 통합 검색으로\n온보딩 시간 대폭 단축",
     PLUM),
]
for i, (title, prob, sol, col) in enumerate(pain_cards):
    x = LM + i * (C3W + Inches(0.28))
    acard(sl, x, cy, C3W, Inches(5.6), accent=col)
    T(sl, title, x + Inches(0.14), cy + Inches(0.1),
      C3W - Inches(0.2), Inches(0.44), sz=14, bold=True, color=col)
    R(sl, x + Inches(0.14), cy + Inches(0.56),
      C3W - Inches(0.28), Inches(0.01), fill=LINE)

    T(sl, "현황 (Pain Point)", x + Inches(0.14), cy + Inches(0.7),
      C3W - Inches(0.2), Inches(0.3), sz=11, bold=True, color=RUST)
    R(sl, x + Inches(0.14), cy + Inches(1.0),
      C3W - Inches(0.28), Inches(1.55),
      fill=RUST_L, lc=BORDER, lw=Pt(0.5))
    T(sl, prob, x + Inches(0.24), cy + Inches(1.1),
      C3W - Inches(0.48), Inches(1.35), sz=13, color=TXT_M)

    T(sl, "MCP 활용 시", x + Inches(0.14), cy + Inches(2.72),
      C3W - Inches(0.2), Inches(0.3), sz=11, bold=True, color=SAGE)
    R(sl, x + Inches(0.14), cy + Inches(3.02),
      C3W - Inches(0.28), Inches(1.55),
      fill=SAGE_L, lc=BORDER, lw=Pt(0.5))
    T(sl, sol, x + Inches(0.24), cy + Inches(3.12),
      C3W - Inches(0.48), Inches(1.35), sz=13, color=TXT_M)

hbar(sl, "OpenAlio MCP — 법령·지침·편람 통합 검색으로 3가지 Pain Point를 단번에 해결",
     accent=NAVY)


# ════════════════════════════════════════════════════════════════
#  S05 — 재경부 공무원 Pain Point
# ════════════════════════════════════════════════════════════════
sl, cy = new_slide(5, "05   재경부·주무부처 공무원 입장에서의 Pain Point",
                   "공개된 데이터를 활용하는 데 막대한 행정비용이 발생")

# 왼쪽: 업무 유형
acard(sl, LM, cy, Inches(3.8), Inches(5.55))
ctitle(sl, "주요 업무 유형", LM, cy + Inches(0.1), Inches(3.8))
R(sl, LM + Inches(0.14), cy + Inches(0.54),
  Inches(3.52), Inches(0.01), fill=LINE)
tasks = ["기능개혁 검토", "증원심사", "경영평가 지원",
         "기관 통합 검토", "보도 대응", "팩트체크"]
ty2 = cy + Inches(0.7)
for t in tasks:
    R(sl, LM + Inches(0.14), ty2, Inches(3.52), Inches(0.55),
      fill=CARD, lc=BORDER, lw=Pt(0.5))
    T(sl, t, LM + Inches(0.28), ty2 + Inches(0.1),
      Inches(3.2), Inches(0.35), sz=14, color=TXT_D, bold=True)
    ty2 += Inches(0.65)

# 오른쪽: 업무 흐름
rx = LM + Inches(4.1)
rw = CW - Inches(4.1)
acard(sl, rx, cy, rw, Inches(5.55), accent=RUST)
ctitle(sl, "실제 업무 흐름 (현재)", rx, cy + Inches(0.1), rw, accent=RUST)
R(sl, rx + Inches(0.14), cy + Inches(0.54),
  rw - Inches(0.28), Inches(0.01), fill=LINE)

flows = [
    ("기관별 정원·재무·사업 정보 확인 필요", WHITE, TXT_D, False),
    ("ALIO 접속 → 항목별 클릭 → 다운로드", RUST_L, RUST, True),
    ("엑셀 파일 수십 개 수작업 정리·가공  ⚠", RUST_L, RUST, True),
    ("분석 수행", WHITE, TXT_M, False),
    ("보고서 작성", WHITE, TXT_M, False),
]
fy2 = cy + Inches(0.7)
for label, bg, col, emph in flows:
    R(sl, rx + Inches(0.14), fy2, rw - Inches(0.28), Inches(0.55),
      fill=bg, lc=BORDER, lw=Pt(0.5))
    if emph:
        R(sl, rx + Inches(0.14), fy2, Inches(0.04), Inches(0.55), fill=RUST)
    T(sl, label, rx + Inches(0.28), fy2 + Inches(0.1),
      rw - Inches(0.5), Inches(0.35), sz=13, bold=emph, color=col)
    if label != "보고서 작성":
        arr_dn(sl, rx + Inches(0.14) + (rw - Inches(0.28)) / 2,
               fy2 + Inches(0.55), Inches(0.18))
    fy2 += Inches(0.73)

T(sl, "⚠  기관 1개당 수십 개 파일  ×  355개 기관 = 막대한 행정비용",
  rx + Inches(0.14), cy + Inches(5.22), rw - Inches(0.28), Inches(0.3),
  sz=11, bold=True, color=RUST, italic=True)

hbar(sl, "OpenAlio MCP — 자동화된 데이터 수집·정리로 행정비용 획기적 절감", accent=NAVY)


# ════════════════════════════════════════════════════════════════
#  S06 — 연구자 Pain Point
# ════════════════════════════════════════════════════════════════
sl, cy = new_slide(6, "06   연구자 입장에서의 Pain Point",
                   "데이터보다 데이터 준비에 더 많은 시간이 소요")

T(sl, "연구 시간의 대부분이 데이터 수집·정제에 소요됩니다",
  LM, cy, CW, Inches(0.42), sz=16, bold=True, color=TXT_D,
  align=PP_ALIGN.CENTER)
cy += Inches(0.52)

cycle = [
    ("수  집", "ALIO·기관 홈페이지\n수작업 다운로드", NAVY),
    ("정  제", "이상값·결측값\n수동 처리", RUST),
    ("통  합", "기관 코드 통일\n형식 표준화", TEAL),
    ("전처리", "분석 가능한\n형태로 변환", PLUM),
]
bw = (CW - Inches(0.28) * 3) / 4
bh = Inches(2.8)
bx = LM
for i, (ctit, cdesc, col) in enumerate(cycle):
    acard(sl, bx, cy, bw, bh, accent=col)
    T(sl, ctit, bx + Inches(0.14), cy + Inches(0.15),
      bw - Inches(0.2), Inches(0.5), sz=17, bold=True, color=col)
    R(sl, bx + Inches(0.14), cy + Inches(0.68),
      bw - Inches(0.28), Inches(0.01), fill=LINE)
    T(sl, cdesc, bx + Inches(0.14), cy + Inches(0.82),
      bw - Inches(0.2), Inches(1.8), sz=13, color=TXT_M)
    if i < 3:
        R(sl, bx + bw, cy + bh / 2 - Inches(0.025),
          Inches(0.28), Inches(0.05), fill=TXT_XL)
    bx += bw + Inches(0.28)

# 반복 루프 표현
R(sl, LM, cy + bh + Inches(0.28), CW, Inches(0.01), fill=LINE)
T(sl, "⟳  위 과정을 매 연구마다, 연구자마다 반복 수행  —  연구 준비기간 수 주",
  LM, cy + bh + Inches(0.34), CW, Inches(0.32),
  sz=13, color=TXT_L, align=PP_ALIGN.CENTER, italic=True)

# 해결책
R(sl, LM, cy + bh + Inches(0.8), CW, Inches(0.82), fill=SAGE_L,
  lc=BORDER, lw=Pt(0.5))
R(sl, LM, cy + bh + Inches(0.8), Inches(0.048), Inches(0.82), fill=SAGE)
T(sl, "OpenAlio MCP 활용 시",
  LM + Inches(0.18), cy + bh + Inches(0.88), Inches(3), Inches(0.35),
  sz=12, bold=True, color=SAGE)
T(sl, "355개 기관 × 11개 메트릭 × 6년치 데이터를 즉시 제공  ·  출처·공시일·단위 자동 첨부  ·  데이터 준비 시간 90% 단축",
  LM + Inches(3.3), cy + bh + Inches(0.88), CW - Inches(3.5), Inches(0.58),
  sz=13, color=TXT_M)


# ════════════════════════════════════════════════════════════════
#  S07 — OpenAlio MCP 개발
# ════════════════════════════════════════════════════════════════
sl, cy = new_slide(7, "07   그래서 OpenAlio MCP를 개발",
                   "공공기관 정보를 AI가 활용할 수 있도록 연결하는 개방형 인터페이스")

# 3레이어 아키텍처
# 레이어 1: 데이터 소스
T(sl, "[ 데이터 소스 레이어 ]", LM, cy, Inches(5), Inches(0.35),
  sz=11, bold=True, color=TXT_L)
cy += Inches(0.38)
src_items2 = [
    ("ALIO", NAVY), ("ALIO PLUS", TEAL), ("JOB-ALIO", PLUM),
    ("법령·지침", SAGE), ("편람", RUST), ("뉴스", GOLD),
]
sw = (CW - Inches(0.2) * 5) / 6
sx2 = LM
for sname, scol in src_items2:
    R(sl, sx2, cy, sw, Inches(0.58), fill=WHITE, lc=scol, lw=Pt(1.2))
    R(sl, sx2, cy, sw, Inches(0.04), fill=scol)
    T(sl, sname, sx2, cy + Inches(0.1), sw, Inches(0.44),
      sz=13, bold=True, color=scol, align=PP_ALIGN.CENTER)
    R(sl, sx2 + sw / 2 - Inches(0.025), cy + Inches(0.58),
      Inches(0.05), Inches(0.28), fill=TXT_XL)
    sx2 += sw + Inches(0.2)
cy += Inches(0.9)

# 레이어 2: MCP (핵심)
R(sl, LM, cy, CW, Inches(1.1), fill=NAVY_D)
R(sl, LM, cy, CW, Inches(0.05), fill=GOLD)
R(sl, LM, cy + Inches(1.05), CW, Inches(0.05), fill=GOLD)
T(sl, "OpenAlio MCP",
  LM, cy + Inches(0.1), CW * 0.5, Inches(0.55),
  sz=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
T(sl, "32개 도구  ·  2개 프롬프트  ·  5개 리소스",
  LM + CW * 0.5, cy + Inches(0.22), CW * 0.5, Inches(0.38),
  sz=14, color=NAVY_L, align=PP_ALIGN.CENTER)
R(sl, LM, cy + Inches(0.58), Inches(0.01), Inches(0.52), fill=GOLD)
T(sl, "표준화된 MCP 프로토콜로 어떤 AI에도 동일한 방식으로 연결",
  LM + CW * 0.5, cy + Inches(0.56), CW * 0.5, Inches(0.38),
  sz=12, color=NAVY_L, italic=True, align=PP_ALIGN.CENTER)

# 화살표
for xi in [LM + CW * 0.1, LM + CW * 0.3, LM + CW * 0.5,
           LM + CW * 0.7, LM + CW * 0.9]:
    R(sl, xi - Inches(0.025), cy + Inches(1.15),
      Inches(0.05), Inches(0.25), fill=TXT_XL)
cy += Inches(1.4)

# 레이어 3: 활용 대상
T(sl, "[ 활용 대상 ]", LM, cy, Inches(5), Inches(0.35),
  sz=11, bold=True, color=TXT_L)
cy += Inches(0.38)
users3 = [
    ("국민", "정보 조회·질문", NAVY),
    ("공공기관", "벤치마킹·학습", TEAL),
    ("공무원", "업무 자동화", PLUM),
    ("연구자", "데이터 분석", SAGE),
]
uw = (CW - Inches(0.35) * 3) / 4
ux2 = LM
for uname, udesc, ucol in users3:
    acard(sl, ux2, cy, uw, Inches(0.88), accent=ucol)
    T(sl, uname, ux2 + Inches(0.14), cy + Inches(0.06),
      uw - Inches(0.2), Inches(0.38), sz=15, bold=True, color=ucol)
    T(sl, udesc, ux2 + Inches(0.14), cy + Inches(0.46),
      uw - Inches(0.2), Inches(0.3), sz=12, color=TXT_L)
    ux2 += uw + Inches(0.35)


# ════════════════════════════════════════════════════════════════
#  S08 — MCP란?
# ════════════════════════════════════════════════════════════════
sl, cy = new_slide(8, "08   MCP란 무엇인가",
                   "Model Context Protocol — AI 시대의 USB-C")

# 왼쪽: USB 비유
acard(sl, LM, cy, C2W, Inches(5.55))
ctitle(sl, "USB-C 비유", LM, cy + Inches(0.1), C2W)
R(sl, LM + Inches(0.14), cy + Inches(0.54),
  C2W - Inches(0.28), Inches(0.01), fill=LINE)

T(sl, "Before  —  기기마다 다른 충전기",
  LM + Inches(0.14), cy + Inches(0.68), C2W - Inches(0.28), Inches(0.32),
  sz=12, bold=True, color=RUST)
devs = ["스마트폰", "태블릿", "노트북", "카메라"]
dx5 = LM + Inches(0.18)
for d in devs:
    R(sl, dx5, cy + Inches(1.05), Inches(1.2), Inches(0.45),
      fill=RUST_L, lc=BORDER, lw=Pt(0.5))
    T(sl, d, dx5, cy + Inches(1.07), Inches(1.2), Inches(0.4),
      sz=11, color=RUST, align=PP_ALIGN.CENTER)
    R(sl, dx5 + Inches(1.2), cy + Inches(1.18),
      Inches(0.22), Inches(0.08), fill=RUST)
    R(sl, dx5 + Inches(1.42), cy + Inches(1.05),
      Inches(0.22), Inches(0.45), fill=RUST)
    dx5 += Inches(1.44)

R(sl, LM + Inches(0.14), cy + Inches(1.7),
  C2W - Inches(0.28), Inches(0.01), fill=LINE)

T(sl, "After  —  USB-C 하나로 통일",
  LM + Inches(0.14), cy + Inches(1.84), C2W - Inches(0.28), Inches(0.32),
  sz=12, bold=True, color=SAGE)
dx6 = LM + Inches(0.18)
ctr = LM + Inches(0.18) + Inches(1.44) * 2 + Inches(0.6)
for d in devs:
    R(sl, dx6, cy + Inches(2.22), Inches(1.2), Inches(0.45),
      fill=SAGE_L, lc=BORDER, lw=Pt(0.5))
    T(sl, d, dx6, cy + Inches(2.24), Inches(1.2), Inches(0.4),
      sz=11, color=SAGE, align=PP_ALIGN.CENTER)
    R(sl, dx6 + Inches(0.57), cy + Inches(2.67),
      Inches(0.05), Inches(0.3), fill=TXT_XL)
    dx6 += Inches(1.44)
R(sl, LM + Inches(0.75), cy + Inches(2.97), C2W - Inches(1.0), Inches(0.05), fill=SAGE)
R(sl, ctr, cy + Inches(3.02), Inches(0.05), Inches(0.3), fill=SAGE)
R(sl, ctr - Inches(0.42), cy + Inches(3.32), Inches(0.9), Inches(0.42), fill=SAGE)
T(sl, "USB-C", ctr - Inches(0.42), cy + Inches(3.34),
  Inches(0.9), Inches(0.38), sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 오른쪽: MCP 설명
acard(sl, C2X2, cy, C2W, Inches(5.55), accent=TEAL)
ctitle(sl, "MCP = Model Context Protocol", C2X2, cy + Inches(0.1), C2W, accent=TEAL)
R(sl, C2X2 + Inches(0.14), cy + Inches(0.54),
  C2W - Inches(0.28), Inches(0.01), fill=LINE)

T(sl, "AI 모델이 다양한 외부 데이터·도구와\n표준화된 방식으로 연결하는 프로토콜",
  C2X2 + Inches(0.14), cy + Inches(0.68), C2W - Inches(0.28), Inches(0.9),
  sz=14, color=TXT_D)

mcp_feats2 = [
    (TEAL,  "표준화", "어떤 AI든 동일한 방식으로 연결"),
    (NAVY,  "안전성", "입력 검증 · 응답 크기 제한 내장"),
    (PLUM,  "모듈성", "도구 단위로 독립 추가·제거"),
    (SAGE,  "확장성", "새 데이터소스 연결 용이"),
    (GOLD,  "투명성", "출처·공시일·단위 자동 첨부"),
]
mfy = cy + Inches(1.72)
for col, mtit, mdesc in mcp_feats2:
    R(sl, C2X2 + Inches(0.14), mfy, C2W - Inches(0.28), Inches(0.6),
      fill=WHITE, lc=BORDER, lw=Pt(0.5))
    R(sl, C2X2 + Inches(0.14), mfy, Inches(0.04), Inches(0.6), fill=col)
    T(sl, mtit, C2X2 + Inches(0.28), mfy + Inches(0.08),
      Inches(1.4), Inches(0.3), sz=13, bold=True, color=col)
    T(sl, mdesc, C2X2 + Inches(1.8), mfy + Inches(0.14),
      C2W - Inches(2.0), Inches(0.32), sz=12, color=TXT_L)
    mfy += Inches(0.7)

T(sl, "Claude · ChatGPT · Cursor · 브리티웍스 등",
  C2X2 + Inches(0.14), cy + Inches(5.22), C2W - Inches(0.28), Inches(0.3),
  sz=11, color=TEAL, italic=True)


# ════════════════════════════════════════════════════════════════
#  S09 — 동작 원리
# ════════════════════════════════════════════════════════════════
sl, cy = new_slide(9, "09   OpenAlio MCP는 어떻게 동작하는가",
                   "질문 한 번으로 공공기관 데이터 분석·보고서까지")

# 중앙 플로우 (좁게)
fw = Inches(5.2)
fx = (W - fw) / 2
fsteps = [
    ("사용자 질문",             NAVY_D),
    ("AI  (Claude / ChatGPT 등)", NAVY),
    ("OpenAlio MCP  (32개 도구)", GOLD),
    ("ALIO · 법령 · 뉴스 조회",  TEAL),
    ("분석 수행",               TEAL),
    ("결과 + 출처 · 공시일 반환", SAGE),
]
fy3 = cy + Inches(0.1)
bh2 = Inches(0.56)
gp  = Inches(0.2)
for i, (lbl, col) in enumerate(fsteps):
    R(sl, fx, fy3, fw, bh2, fill=col)
    T(sl, lbl, fx + Inches(0.1), fy3 + Inches(0.08),
      fw - Inches(0.2), bh2 - Inches(0.14), sz=14, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(fsteps) - 1:
        arr_dn(sl, fx + fw / 2, fy3 + bh2, gp)
    fy3 += bh2 + gp

# 왼쪽: 예시 질문
acard(sl, LM, cy, Inches(3.3), Inches(5.5))
ctitle(sl, "예시 질문", LM, cy + Inches(0.1), Inches(3.3))
R(sl, LM + Inches(0.14), cy + Inches(0.54),
  Inches(2.92), Inches(0.01), fill=LINE)
qs2 = [
    '"최근 5년 부채 증가율\n상위 10개 기관"',
    '"A기관·B기관 통합\n검토 보고서"',
    '"육아휴직 사용률 가장\n높은 기관은?"',
]
qy3 = cy + Inches(0.68)
for q in qs2:
    R(sl, LM + Inches(0.14), qy3, Inches(2.92), Inches(1.3),
      fill=NAVY_L, lc=BORDER, lw=Pt(0.5))
    T(sl, q, LM + Inches(0.24), qy3 + Inches(0.16),
      Inches(2.72), Inches(1.0), sz=13, color=NAVY_D, italic=True)
    qy3 += Inches(1.48)

# 오른쪽: 출력 형태
acard(sl, W - LM - Inches(3.3), cy, Inches(3.3), Inches(5.5), accent=SAGE)
ctitle(sl, "출력 형태", W - LM - Inches(3.3), cy + Inches(0.1),
       Inches(3.3), accent=SAGE)
R(sl, W - LM - Inches(3.16), cy + Inches(0.54),
  Inches(2.92), Inches(0.01), fill=LINE)
outputs3 = [
    ("시계열 차트", "연도별 추이 자동 시각화"),
    ("비교 표", "기관 간 지표 나란히 비교"),
    ("보고서 초안", "분석 결과 정형화된 텍스트"),
    ("출처 첨부", "시스템명·공시일·단위 포함"),
    ("결측 안내", "미공시·공시주기 안내 포함"),
]
oy5 = cy + Inches(0.68)
ox5 = W - LM - Inches(3.16)
for otit, odesc in outputs3:
    R(sl, ox5, oy5, Inches(2.92), Inches(0.82),
      fill=SAGE_L, lc=BORDER, lw=Pt(0.5))
    T(sl, otit, ox5 + Inches(0.1), oy5 + Inches(0.06),
      Inches(2.7), Inches(0.32), sz=13, bold=True, color=SAGE)
    T(sl, odesc, ox5 + Inches(0.1), oy5 + Inches(0.42),
      Inches(2.7), Inches(0.32), sz=11, color=TXT_L)
    oy5 += Inches(0.95)


# ════════════════════════════════════════════════════════════════
#  S10 — 활용 사례 ① 국민
# ════════════════════════════════════════════════════════════════
sl, cy = new_slide(10, "10   활용 사례 ①  국민 서비스",
                   "누구나 AI에게 물어보면 공공기관 정보를 즉시 확인")

cases_c = [
    ("한전 부채 얼마야?",
     "한국전력 최근 5년 부채 추이\n수치·그래프 즉시 제공", NAVY),
    ("육아휴직 사용률이\n가장 높은 기관은?",
     "355개 기관 육아휴직 데이터\n순위 표 자동 생성", TEAL),
    ("근처 공공기관\n체육시설 알려줘",
     "ALIO PLUS 시설 정보\n예약 링크·운영시간 안내", PLUM),
]
for i, (q, res, col) in enumerate(cases_c):
    x = LM + i * (C3W + Inches(0.28))
    acard(sl, x, cy, C3W, Inches(5.55), accent=col)
    # 질문
    R(sl, x + Inches(0.14), cy + Inches(0.15),
      C3W - Inches(0.28), Inches(1.35), fill=CARD, lc=BORDER, lw=Pt(0.5))
    T(sl, "질문", x + Inches(0.24), cy + Inches(0.22),
      C3W - Inches(0.48), Inches(0.28), sz=10, bold=True, color=TXT_L)
    T(sl, f'"{q}"', x + Inches(0.24), cy + Inches(0.52),
      C3W - Inches(0.48), Inches(0.9), sz=14, bold=True, color=col)
    # 결과
    R(sl, x + Inches(0.14), cy + Inches(1.65),
      C3W - Inches(0.28), Inches(0.01), fill=LINE)
    T(sl, "→  결과", x + Inches(0.14), cy + Inches(1.8),
      C3W - Inches(0.28), Inches(0.3), sz=11, bold=True, color=col)
    R(sl, x + Inches(0.14), cy + Inches(2.16),
      C3W - Inches(0.28), Inches(2.72), fill=WHITE, lc=BORDER, lw=Pt(0.5))
    T(sl, res, x + Inches(0.24), cy + Inches(2.3),
      C3W - Inches(0.48), Inches(2.5), sz=14, color=TXT_M)

hbar(sl, "사용자는 웹사이트 구조를 몰라도 됩니다  —  질문만 하면 됩니다", accent=TEAL)


# ════════════════════════════════════════════════════════════════
#  S11 — 활용 사례 ② 공공기관 실무자
# ════════════════════════════════════════════════════════════════
sl, cy = new_slide(11, "11   활용 사례 ②  공공기관 실무자",
                   "벤치마킹·평가 대응·조직운영 정보를 즉시 활용")

cases_i = [
    ("유사기관 ESG 우수사례 정리",
     '"유사기관 ESG 우수사례 정리"',
     "동종 기관 ESG 공시 항목\n우수 사례 자동 요약 제공", NAVY),
    ("경영평가 지적사항 분석",
     '"최근 경영평가 지적사항 분석"',
     "2025·2026 평가편람 기준\n지적 패턴·주요 지표 분석", TEAL),
    ("기관별 조직운영 비교",
     '"기관별 조직운영 비교"',
     "정원·실인원·보수 등\n11개 메트릭 동시 비교", PLUM),
]
for i, (title, q, res, col) in enumerate(cases_i):
    x = LM + i * (C3W + Inches(0.28))
    acard(sl, x, cy, C3W, Inches(5.55), accent=col)
    ctitle(sl, title, x, cy + Inches(0.1), C3W, accent=col)
    R(sl, x + Inches(0.14), cy + Inches(0.54),
      C3W - Inches(0.28), Inches(0.01), fill=LINE)
    R(sl, x + Inches(0.14), cy + Inches(0.68),
      C3W - Inches(0.28), Inches(0.9), fill=CARD, lc=BORDER, lw=Pt(0.5))
    T(sl, "질문", x + Inches(0.24), cy + Inches(0.74),
      C3W - Inches(0.48), Inches(0.26), sz=10, bold=True, color=TXT_L)
    T(sl, q, x + Inches(0.24), cy + Inches(1.0),
      C3W - Inches(0.48), Inches(0.5), sz=13, color=col, italic=True)
    R(sl, x + Inches(0.14), cy + Inches(1.74),
      C3W - Inches(0.28), Inches(0.01), fill=LINE)
    T(sl, "출력 결과", x + Inches(0.14), cy + Inches(1.88),
      C3W - Inches(0.28), Inches(0.3), sz=11, bold=True, color=col)
    R(sl, x + Inches(0.14), cy + Inches(2.24),
      C3W - Inches(0.28), Inches(2.65), fill=WHITE, lc=BORDER, lw=Pt(0.5))
    T(sl, res, x + Inches(0.24), cy + Inches(2.38),
      C3W - Inches(0.48), Inches(2.45), sz=14, color=TXT_M)

hbar(sl, "법령·지침·편람 통합 검색  —  신규 담당자 온보딩 시간 대폭 단축", accent=TEAL)


# ════════════════════════════════════════════════════════════════
#  S12 — 활용 사례 ③ 재경부 (★ 핵심)
# ════════════════════════════════════════════════════════════════
sl, cy = new_slide(12, "12   활용 사례 ③  재경부 업무  ★",
                   "기관 통합 검토 보고서를 AI가 자동으로 작성")

# 질문 강조
R(sl, LM, cy, CW, Inches(0.66), fill=NAVY_D)
R(sl, LM, cy, Inches(0.055), Inches(0.66), fill=GOLD)
T(sl, '  "A기관과 B기관의 통합 검토 보고서를 작성해 줘"',
  LM + Inches(0.18), cy + Inches(0.11), CW - Inches(0.38), Inches(0.46),
  sz=18, bold=True, color=WHITE)
cy += Inches(0.78)

# 출력 6항목 (3+3 그리드)
outw = (CW - Inches(0.22) * 2) / 3 * 0.9
outitems = [
    ("재무현황", "부채·자산·수익 최근 5년 추이", NAVY),
    ("인력현황", "정원·실인원·신규 채용 현황", TEAL),
    ("주요사업", "설립 목적·핵심 기능 정리", PLUM),
    ("지역분포", "본사·지사 소재지·지역별 인력", RUST),
    ("중복기능", "유사 기능 도출·통합 시 효율", SAGE),
    ("기대효과", "시나리오별 예상 효과 요약", GOLD),
]
ow = (CW * 0.6 - Inches(0.2)) / 3
for i, (otit, odesc, ocol) in enumerate(outitems):
    ox6 = LM + (i % 3) * (ow + Inches(0.1))
    oy6 = cy + (i // 3) * (Inches(1.2) + Inches(0.15))
    acard(sl, ox6, oy6, ow, Inches(1.2), accent=ocol)
    T(sl, otit, ox6 + Inches(0.14), oy6 + Inches(0.08),
      ow - Inches(0.2), Inches(0.4), sz=14, bold=True, color=ocol)
    T(sl, odesc, ox6 + Inches(0.14), oy6 + Inches(0.52),
      ow - Inches(0.2), Inches(0.6), sz=11, color=TXT_L)

# 오른쪽: 기대효과
ex = LM + CW * 0.62
ew = CW * 0.38
acard(sl, ex, cy, ew, Inches(2.58), accent=GOLD)
ctitle(sl, "기대 효과", ex, cy + Inches(0.1), ew, accent=GOLD)
R(sl, ex + Inches(0.14), cy + Inches(0.54),
  ew - Inches(0.28), Inches(0.01), fill=LINE)
effects3 = [
    "수 시간 → 수 분 보고서 작성",
    "수작업 엑셀 없이 자동 수집",
    "출처 자동 첨부 (공시일 포함)",
    "다양한 업무에 재사용 가능",
]
efy2 = cy + Inches(0.68)
for eff in effects3:
    R(sl, ex + Inches(0.14), efy2,
      ew - Inches(0.28), Inches(0.44), fill=GOLD_L, lc=BORDER, lw=Pt(0.5))
    R(sl, ex + Inches(0.14), efy2,
      Inches(0.04), Inches(0.44), fill=GOLD)
    T(sl, eff, ex + Inches(0.26), efy2 + Inches(0.08),
      ew - Inches(0.44), Inches(0.3), sz=12, color=TXT_D)
    efy2 += Inches(0.52)

# 브리티웍스/내부시스템 연결 힌트
acard(sl, ex, cy + Inches(2.76), ew, Inches(0.6), accent=TEAL)
T(sl, "브리티웍스·지능형 업무시스템에\nMCP 연결 시 즉시 활용 가능",
  ex + Inches(0.14), cy + Inches(2.82), ew - Inches(0.2), Inches(0.5),
  sz=12, color=TXT_M)

hbar(sl, "공개된 데이터  +  AI  =  행정 생산성의 획기적 향상", accent=GOLD)


# ════════════════════════════════════════════════════════════════
#  S13 — 브리티웍스 연계
# ════════════════════════════════════════════════════════════════
sl, cy = new_slide(13, "13   활용 사례 ④  브리티웍스 및 지능형 업무시스템",
                   "OpenAlio MCP는 특정 AI 서비스가 아닌 공공기관 데이터 기반 인프라")

T(sl, "표준화된 MCP 프로토콜로 어떤 AI 시스템에도 연결 가능한 개방형 인터페이스",
  LM, cy, CW, Inches(0.38), sz=15, bold=True, color=TXT_D, align=PP_ALIGN.CENTER)
cy += Inches(0.5)

# AI 서비스 레이어
ai_svcs2 = [
    ("브리티웍스", "공공기관 내부 업무", NAVY),
    ("지능형 업무시스템", "행정 AI 플랫폼", TEAL),
    ("Claude / ChatGPT", "범용 AI 어시스턴트", PLUM),
    ("향후 AI 플랫폼", "미래 공공 AI 서비스", SAGE),
]
aw2 = (CW - Inches(0.3) * 3) / 4
ax4 = LM
for aname, adesc, acol in ai_svcs2:
    acard(sl, ax4, cy, aw2, Inches(0.92), accent=acol)
    T(sl, aname, ax4 + Inches(0.14), cy + Inches(0.08),
      aw2 - Inches(0.2), Inches(0.38), sz=13, bold=True, color=acol)
    T(sl, adesc, ax4 + Inches(0.14), cy + Inches(0.5),
      aw2 - Inches(0.2), Inches(0.3), sz=11, color=TXT_L)
    R(sl, ax4 + aw2 / 2 - Inches(0.025), cy + Inches(0.92),
      Inches(0.05), Inches(0.32), fill=TXT_XL)
    ax4 += aw2 + Inches(0.3)
cy += Inches(1.28)

# MCP 핵심 박스
R(sl, LM, cy, CW, Inches(0.95), fill=NAVY_D)
R(sl, LM, cy, CW, Inches(0.04), fill=GOLD)
R(sl, LM, cy + Inches(0.91), CW, Inches(0.04), fill=GOLD)
T(sl, "OpenAlio MCP  —  표준 인터페이스",
  LM, cy + Inches(0.1), CW * 0.6, Inches(0.52),
  sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
T(sl, "표준화된 MCP 프로토콜 · 어떤 AI든 동일하게 연결",
  LM + CW * 0.6, cy + Inches(0.24), CW * 0.4, Inches(0.35),
  sz=13, color=NAVY_L, italic=True, align=PP_ALIGN.CENTER)

# 화살표
for xi2 in [LM + CW * 0.15, LM + CW * 0.38,
            LM + CW * 0.62, LM + CW * 0.85]:
    R(sl, xi2, cy + Inches(0.95), Inches(0.05), Inches(0.3), fill=TXT_XL)
cy += Inches(1.29)

# 데이터 소스 레이어
dw2 = (CW - Inches(0.28) * 3) / 4
dx7 = LM
for dname, dcol in [("ALIO", NAVY), ("ALIO PLUS", TEAL),
                    ("JOB-ALIO", PLUM), ("법령·지침·편람", SAGE)]:
    R(sl, dx7, cy, dw2, Inches(0.62), fill=WHITE, lc=dcol, lw=Pt(1.2))
    R(sl, dx7, cy, dw2, Inches(0.038), fill=dcol)
    T(sl, dname, dx7, cy + Inches(0.1), dw2, Inches(0.44),
      sz=13, bold=True, color=dcol, align=PP_ALIGN.CENTER)
    dx7 += dw2 + Inches(0.28)
cy += Inches(0.82)

# 설명 포인트
R(sl, LM, cy, CW, Inches(1.25), fill=WHITE, lc=BORDER, lw=Pt(0.5))
R(sl, LM, cy, Inches(0.048), Inches(1.25), fill=TEAL)
pts3 = [
    "표준화된 MCP 프로토콜로 브리티웍스·지능형 업무시스템에 바로 연결 가능",
    "새 AI 플랫폼 도입 시에도 OpenAlio MCP는 그대로 재사용  —  인프라 투자 보호",
    "데이터 추가(평가결과·국회·조달)도 MCP 서버에만 반영하면 모든 AI에 즉시 적용",
]
py4 = cy + Inches(0.1)
for pt in pts3:
    T(sl, f"·  {pt}", LM + Inches(0.18), py4,
      CW - Inches(0.3), Inches(0.34), sz=13, color=TXT_M)
    py4 += Inches(0.36)


# ════════════════════════════════════════════════════════════════
#  S14 — 향후 발전 방향
# ════════════════════════════════════════════════════════════════
sl, cy = new_slide(14, "14   향후 발전 방향",
                   "단계적 확장을 통해 공공기관 정보공개 체계를 AI 시대에 맞게 고도화")

roadmap3 = [
    ("대국민 서비스", NAVY, [
        "ALIO 챗봇 구현",
        "기관 비교 서비스",
        "정책 팩트체크",
    ]),
    ("내부 업무 지원", TEAL, [
        "증원심사 지원",
        "기능개혁 지원",
        "경영평가 지원",
        "언론대응 · 기관현황 작성",
    ]),
    ("연구 활용", PLUM, [
        "공공기관 데이터 분석 플랫폼",
        "평가결과 데이터 추가",
        "국회 · 조달 데이터 연계",
    ]),
    ("인프라 고도화", GOLD, [
        "HTTP/SSE 서버 전환",
        "인증 · 속도 제한 적용",
        "자동 업데이트 파이프라인",
        "보안 심의 대응",
    ]),
]
rw2 = C3W  # 4 cols same as 3col but 4
rw2 = (CW - Inches(0.28) * 3) / 4
rx3 = LM
for rname, rcol, ritems in roadmap3:
    acard(sl, rx3, cy, rw2, Inches(5.55), accent=rcol)
    T(sl, rname, rx3 + Inches(0.14), cy + Inches(0.1),
      rw2 - Inches(0.2), Inches(0.44), sz=14, bold=True, color=rcol)
    R(sl, rx3 + Inches(0.14), cy + Inches(0.56),
      rw2 - Inches(0.28), Inches(0.01), fill=LINE)
    iy3 = cy + Inches(0.7)
    for item in ritems:
        R(sl, rx3 + Inches(0.18), iy3 + Inches(0.12),
          Inches(0.07), Inches(0.07), fill=rcol)
        T(sl, item, rx3 + Inches(0.34), iy3,
          rw2 - Inches(0.48), Inches(0.52), sz=13, color=TXT_M)
        iy3 += Inches(0.62)
    rx3 += rw2 + Inches(0.28)

# 로드맵 타임라인
R(sl, LM, cy + Inches(5.7), CW, Inches(0.5), fill=NAVY_D)
phases2 = ["Phase 1  ✓  완료", "Phase 2  ●  현재", "Phase 3  ○  계획", "Phase 4  ○  장기"]
px2 = LM + Inches(0.5)
pw2 = (CW - Inches(1.0)) / 4
for i, ph in enumerate(phases2):
    T(sl, ph, px2, cy + Inches(5.76), pw2, Inches(0.4),
      sz=12, bold=(i == 1), color=GOLD if i == 1 else NAVY_L,
      align=PP_ALIGN.CENTER)
    px2 += pw2


# ════════════════════════════════════════════════════════════════
#  S15 — 마무리
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, Inches(7.1), H, fill=NAVY_D)
R(sl, Inches(7.1), 0, W - Inches(7.1), H, fill=BG)
# 오른쪽 장식
R(sl, Inches(7.4), Inches(0.5), W - Inches(7.8), Inches(0.01), fill=LINE)
R(sl, Inches(7.4), Inches(1.5), W - Inches(7.8), Inches(0.01), fill=LINE)
R(sl, Inches(7.4), Inches(2.5), W - Inches(7.8), Inches(0.01), fill=LINE)
R(sl, Inches(7.6), Inches(0.5), Inches(0.01), Inches(6.5), fill=LINE)
R(sl, Inches(7.6), Inches(0.5), Inches(0.22), Inches(0.22), fill=NAVY_L)

# 왼쪽 콘텐츠
R(sl, 0, 0, Inches(7.1), Inches(0.055), fill=TEAL)
T(sl, "마무리", Inches(0.55), Inches(0.55), Inches(6.0), Inches(0.42),
  sz=14, color=TEAL_L, italic=True)
T(sl, "정보공개의\n다음 단계는\n정보 활용",
  Inches(0.55), Inches(1.05), Inches(6.1), Inches(1.95),
  sz=38, bold=True, color=WHITE)
R(sl, Inches(0.55), Inches(3.1), Inches(4.0), Inches(0.05), fill=TEAL)

msgs2 = [
    "OpenAlio MCP는 공공기관 정보공개 체계를",
    "AI 시대에 맞게 확장하기 위한 실험입니다.",
    " ",
    "국민과 행정이 공공기관 정보를 더 쉽고",
    "효과적으로 활용할 수 있도록 지원합니다.",
]
my3 = Inches(3.28)
for msg in msgs2:
    T(sl, msg, Inches(0.55), my3, Inches(6.1), Inches(0.4),
      sz=15, color=NAVY_L)
    my3 += Inches(0.4)

# 핵심 키워드 3개
kws3 = [
    ("개방성", "오픈소스 · 누구나 기여", TEAL),
    ("연결성", "AI-데이터 표준 연결", NAVY),
    ("활용성", "질문 한 번 · 즉시 분석", GOLD),
]
kx3 = Inches(0.55)
ky3 = Inches(5.55)
for kword, kdesc, kcol in kws3:
    R(sl, kx3, ky3, Inches(1.85), Inches(1.35), fill=WHITE,
      lc=kcol, lw=Pt(1.5))
    R(sl, kx3, ky3, Inches(1.85), Inches(0.04), fill=kcol)
    T(sl, kword, kx3 + Inches(0.1), ky3 + Inches(0.12),
      Inches(1.65), Inches(0.5), sz=16, bold=True, color=WHITE)
    T(sl, kdesc, kx3 + Inches(0.1), ky3 + Inches(0.65),
      Inches(1.65), Inches(0.55), sz=11, color=NAVY_L)
    kx3 += Inches(2.0)

T(sl, "감사합니다",
  Inches(0.55), Inches(7.0), Inches(4), Inches(0.45),
  sz=26, bold=True, color=GOLD)

# 오른쪽 패널 내용
acard(sl, Inches(7.9), Inches(1.5), Inches(5.0), Inches(4.0), accent=TEAL, bg=WHITE)
ctitle(sl, "핵심 요약", Inches(7.9), Inches(1.6), Inches(5.0), accent=TEAL)
R(sl, Inches(8.04), Inches(2.04), Inches(4.72), Inches(0.01), fill=LINE)
summaries = [
    (NAVY,  "공공기관 데이터", "355개 기관 · 11메트릭 · 6년치"),
    (TEAL,  "AI 도구",       "32개 도구 · 2 프롬프트 · 5 리소스"),
    (PLUM,  "연결 방식",     "MCP 표준 프로토콜 (개방형)"),
    (GOLD,  "활용 대상",     "국민 · 공무원 · 연구자"),
    (SAGE,  "현재 상태",     "Phase 2 — 로컬 MCP 서버"),
    (RUST,  "다음 단계",     "Phase 4 — HTTP 서버 · 인증 · 자동화"),
]
sy5 = Inches(2.18)
for scol, stit, sdesc in summaries:
    R(sl, Inches(8.04), sy5, Inches(4.72), Inches(0.52),
      fill=WHITE, lc=BORDER, lw=Pt(0.5))
    R(sl, Inches(8.04), sy5, Inches(0.04), Inches(0.52), fill=scol)
    T(sl, stit, Inches(8.22), sy5 + Inches(0.06),
      Inches(1.4), Inches(0.3), sz=12, bold=True, color=scol)
    T(sl, sdesc, Inches(9.7), sy5 + Inches(0.12),
      Inches(2.9), Inches(0.3), sz=12, color=TXT_L)
    sy5 += Inches(0.58)


# ── 저장 ──────────────────────────────────────────────────────
out = "/home/user/open-ALIO-mcp/OpenAlio_MCP_발표자료.pptx"
prs.save(out)
print(f"저장 완료: {out}")
